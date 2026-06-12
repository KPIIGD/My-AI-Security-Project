"""
PPTX 재구성 — 본편 11장 + Appendix 구분자 + 부록 17장.
슬라이드 내용은 불변. 순서만 바꾸고 구분자 한 장 추가.
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r"c:\litellm\capstone_pii_presentation_polished.pptx"
DST = r"c:\litellm\capstone_pii_presentation_core.pptx"

CORE = [1, 2, 15, 19, 16, 17, 18, 20, 23, 27, 28]
APPENDIX = [21, 22, 24, 25, 26, 3, 9, 10, 11, 14, 4, 6, 7, 8, 5, 12, 13]

prs = Presentation(SRC)
assert len(prs.slides) == 28

# 1) Add Appendix separator slide at the end (index 28)
layout = prs.slide_layouts[0]
sep = prs.slides.add_slide(layout)
# Strip auto-inherited placeholders so the separator starts clean
for ph in list(sep.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

sw, sh = prs.slide_width, prs.slide_height
tb = sep.shapes.add_textbox(Emu(0), Emu(sh // 2 - 600000), sw, Emu(1200000))
tf = tb.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
r1 = p1.add_run()
r1.text = "Appendix"
r1.font.name = "맑은 고딕"
r1.font.size = Pt(54)
r1.font.bold = True
r1.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "상세 데이터 · 통계 검정 · 방법론 보완"
r2.font.name = "맑은 고딕"
r2.font.size = Pt(18)
r2.font.color.rgb = RGBColor(0x5C, 0x5C, 0x5C)

# 2) Reorder: CORE → separator → APPENDIX
sep_idx = 28  # separator is the 29th slide (0-indexed 28)
new_order_idx = [i - 1 for i in CORE] + [sep_idx] + [i - 1 for i in APPENDIX]
assert sorted(new_order_idx) == list(range(29))

xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
assert len(slides) == 29
new_slides = [slides[i] for i in new_order_idx]
for s in slides:
    xml_slides.remove(s)
for s in new_slides:
    xml_slides.append(s)

prs.save(DST)

import os
mb = os.path.getsize(DST) / 1024 / 1024
print(f"[OK] slides: {len(prs.slides)}")
print(f"[OK] DST: {mb:.1f} MB -> {DST}")
print(f"[OK] Core (1~11): 표지/목차 + 문제 + 결과 + 설계 + 방법 + 4way + Cascade + 결론 + Q&A")
print(f"[OK] Appendix (13~29): McNemar, Ablation, Robustness, 논의, 한계, 배경 단편들")
