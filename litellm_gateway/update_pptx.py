"""Append result slides (with figures) to the existing presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import copy

SRC = "AI_Gateway_가드레일_프로젝트_발표.pptx"
DST = "AI_Gateway_가드레일_프로젝트_발표_v3_results.pptx"

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[0]  # only one layout in this template

NAVY = RGBColor(0x1F, 0x2A, 0x44)
RED = RGBColor(0xD9, 0x53, 0x4F)
GREEN = RGBColor(0x5C, 0xB8, 0x5C)
GRAY = RGBColor(0x66, 0x66, 0x66)


def add_title_slide(title_text, subtitle_text):
    s = prs.slides.add_slide(BLANK)
    box = s.shapes.add_textbox(Inches(0.6), Inches(0.5), SW - Inches(1.2), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle_text:
        sub = s.shapes.add_textbox(Inches(0.6), Inches(1.4), SW - Inches(1.2), Inches(0.6))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle_text
        sp.font.size = Pt(14)
        sp.font.color.rgb = GRAY
    return s


def add_image_slide(title, subtitle, img_path, img_top=Inches(1.4), img_height=None, img_width=None):
    s = add_title_slide(title, subtitle)
    if img_height:
        s.shapes.add_picture(img_path, Inches(0.6), img_top, height=img_height)
    elif img_width:
        s.shapes.add_picture(img_path, Inches(0.6), img_top, width=img_width)
    else:
        s.shapes.add_picture(img_path, Inches(0.6), img_top, width=SW - Inches(1.2))
    return s


def add_table_slide(title, subtitle, headers, rows, col_widths=None):
    s = add_title_slide(title, subtitle)
    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(0.6), Inches(1.6), SW - Inches(1.2), Inches(0.4) * n_rows
    tbl_shape = s.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(12)
                r.font.color.rgb = NAVY
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
    return s


# ── Section divider ──
add_title_slide("실험 결과 (10,000건 진짜 평가)", "2026-04-20  |  LiteLLM Gateway 4계층 호출  |  TRUE detection 기준")

# ── Slide R1: 헤드라인 결론 ──
s = add_title_slide("핵심 결론", "프로덕션 가드레일은 한국어 텍스트형 PII에 무방비, Layer 0가 메운다")
box = s.shapes.add_textbox(Inches(0.6), Inches(1.6), SW - Inches(1.2), Inches(5))
tf = box.text_frame
tf.word_wrap = True
points = [
    ("Baseline 한국어 텍스트형 PII", "47.93% 탐지 (52.07% 우회 — 절반 이상 뚫림)"),
    ("Layer 0 추가 후", "96.08% 탐지 (3.92% 우회) — +48.15%p 회복"),
    ("영어 PII", "99.23% 유지 — Layer 0 영향 없음 (false positive 0건)"),
    ("L0 단독 차단 (다른 가드레일 다 뚫린 케이스)", "1,444건"),
    ("100% 차단된 PII 타입 (11개)", "allergy, company, job_title, gps, degree, blood, marital, dept, jwt, diagnosis, religion"),
]
for i, (k, v) in enumerate(points):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"• {k}: "
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    r = p.add_run()
    r.text = v
    r.font.size = Pt(16)
    r.font.bold = False
    r.font.color.rgb = GRAY

# ── Slide R2: 전체 우회율 figure ──
add_image_slide("Fig 1 — 전체 우회율: Baseline vs Layer 0",
                "한국어 텍스트형 PII에서 격차 가장 큼 (52.1% → 3.9%)",
                "fig1_overall_bypass.png", img_top=Inches(1.6), img_width=SW - Inches(1.2))

# ── Slide R3: Lang × Validity figure ──
add_image_slide("Fig 2 — Language × Validity TRUE detection rate",
                "KR_semantic +48.1%p, KR_format +17.9%p — 한국어 슬라이스에서 회복",
                "fig2_lang_x_validity.png", img_top=Inches(1.6), img_width=SW - Inches(1.2))

# ── Slide R4: Hardest PII figure ──
add_image_slide("Fig 3 — Top 20 hardest PII bypass before/after",
                "11개 PII 타입은 100% 차단 (0% 우회) — 알레르기·회사·직책·학력 등",
                "fig3_hardest_pii.png", img_top=Inches(1.4), img_height=Inches(5.6))

# ── Slide R5: Mutation level figure ──
add_image_slide("Fig 4 — 변이 레벨별 우회율",
                "L4 Linguistic·L5 Context 같은 어려운 변이도 큰 폭 회복",
                "fig4_mutation_level.png", img_top=Inches(1.6), img_width=SW - Inches(1.2))

# ── Slide R6: L0 solo catches figure ──
add_image_slide("Fig 5 — Layer 0 단독 차단 PII 타입",
                "1,444건 중 top 20 — 한국어 텍스트형 PII가 압도적",
                "fig5_l0_solo_pii.png", img_top=Inches(1.4), img_height=Inches(5.6))

# ── Slide R7: 데이터 신뢰성 ──
add_table_slide(
    "데이터 신뢰성",
    "오염 없음 / cherry-pick 없음 / 동일 입력으로 1:1 비교",
    ["항목", "결과"],
    [
        ["평가 케이스 수", "10,000건 (stratified sampling, KR 6,513 / EN 3,487)"],
        ["호출 방식", "LiteLLM Gateway /apply_guardrail (진짜 외부 API 호출)"],
        ["Bedrock errors", "0건 (config 스키마 패치 후)"],
        ["Presidio errors", "0건"],
        ["Lakera errors", "0건 (catch 0%은 정상 — 인젝션 전용)"],
        ["Layer 0 errors", "0건"],
        ["Layer 0 영어 false positive", "0건 (영어에 BLOCK 안 함)"],
        ["Baseline 4/16 vs 4/19 일치", "TRUE 78.5% vs 79.01% — 매우 일관됨"],
        ["슬라이스 정의", "validity_group은 퍼저 단계 a priori 분류 (post-hoc 아님)"],
    ],
    col_widths=[2.0, 6.5],
)

# ── Slide R8: 남은 한계 + 향후 과제 ──
s = add_title_slide("한계 및 향후 과제", "남은 6.27% 우회 (602건)에 대한 대응 방향")
box = s.shapes.add_textbox(Inches(0.6), Inches(1.6), SW - Inches(1.2), Inches(5))
tf = box.text_frame
tf.word_wrap = True
items = [
    "남은 우회 가장 많은 PII 타입: emp_id 39, face_id 35, car_ins 30, EN_SSN 27, rrn 25",
    "Validity group: format 443건 (특수 형식 PII), checksum 111, semantic 48 (L0 거의 다 잡음)",
    "Mutation level: L2 Encoding 235건 (정규화로도 못 잡힌 인코딩 변이)",
    "",
    "향후 과제:",
    "  1. Layer 0 키워드 사전 확장 (emp_id, face_id 패턴)",
    "  2. L2 인코딩 정규화 강화 (combining mark, soft hyphen 등)",
    "  3. RAG 삽입 공격 (ctx_rag) 별도 탐지 — 문서 내 한국어 키워드 매칭",
    "  4. 운영 환경 latency 측정 (현재 L0 추가 latency ~10ms)",
]
for i, t in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = t
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY if "향후" in t else GRAY

prs.save(DST)
print(f"saved: {DST}")
print(f"slides: {len(prs.slides)}")
