# -*- coding: utf-8 -*-
"""
프로그램 아키텍처 — M1~M9 파이프라인만, 슬라이드 정중앙 정렬 (사이드바 없음)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY   = RGBColor(0x16, 0x29, 0x47)
INK    = RGBColor(0x22, 0x2A, 0x35)
SUB    = RGBColor(0x5A, 0x66, 0x77)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PAPER  = RGBColor(0xF4, 0xF6, 0xF9)

C_IN   = RGBColor(0x37, 0x4A, 0x67)
C_M1   = RGBColor(0xE8, 0x6F, 0x1A)
C_M1BG = RGBColor(0xFD, 0xF0, 0xE3)
C_DET  = RGBColor(0x1F, 0x77, 0xB4)
C_DETBG= RGBColor(0xEA, 0xF3, 0xFB)
C_REF  = RGBColor(0x16, 0x8F, 0x7E)
C_REFBG= RGBColor(0xE7, 0xF5, 0xF2)
C_POL  = RGBColor(0x7A, 0x3E, 0xA6)
C_POLBG= RGBColor(0xF1, 0xEA, 0xF8)
C_AUD  = RGBColor(0x99, 0x6A, 0x12)
C_AUDBG= RGBColor(0xFA, 0xF2, 0xDE)

FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = PAPER; bg.line.fill.background(); bg.shadow.inherit = False


def _set(tf, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(7); tf.margin_right = Pt(7); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        sp = spec[4] if len(spec) > 4 else 1
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp); p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
        rPr = r._r.get_or_add_rPr(); rPr.append(rPr.makeelement(qn('a:ea'), {'typeface': FONT}))


def box(x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=1.0, radius=0.10):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: s.adjustments[0] = radius
        except Exception: pass
    return s


def tb(x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set(t.text_frame, lines, align, anchor); return t


def lbox(x, y, w, h, fill, lines, align=PP_ALIGN.LEFT, line=None, line_w=1.0, radius=0.10):
    box(x, y, w, h, fill, line=line, line_w=line_w, radius=radius)
    tb(x, y, w, h, lines, align=align, anchor=MSO_ANCHOR.MIDDLE)


def arrow(cx, y, h=0.20, color=INK, w=0.40):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(cx-w/2), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background(); s.shadow.inherit = False


# ===== 정중앙 컬럼 좌표 =====
LW = 9.60
LX = (13.333 - LW) / 2          # 좌우 정중앙
cx = LX + LW/2                  # = 6.667 (슬라이드 한가운데)

# ===== 타이틀 =====
box(LX, 0.22, LW, 0.84, NAVY, radius=0.14)
tb(LX, 0.22, LW, 0.84, [("프로그램 아키텍처", 26, True, WHITE)], align=PP_ALIGN.CENTER)

# 입력
y = 1.22
lbox(LX, y, LW, 0.52, C_IN,
     [("입력  GuardrailRequest", 13.5, True, WHITE),
      ("text · policy_profile · output_target(llm_input/external/audit) · scan_stage(input/output)", 10, False, RGBColor(0xCF,0xDB,0xEC))],
     radius=0.16)
arrow(cx, y+0.54, color=C_M1)

# M1
y = 1.90
box(LX, y, LW, 0.96, C_M1BG, line=C_M1, line_w=2.0, radius=0.08)
tb(LX, y, LW, 0.96, [
    ("M1  preprocess.py  —  정규화 + 원문 offset map  ⭐", 14, True, C_M1),
    ("투명문자 제거 · 전각/특수숫자→ASCII · 자모/초성체/야민정음/로마자/한글숫자 복원 · 띄어쓰기 압축", 10.5, False, INK),
    ("→ 변이(variant) 9종 생성  +  모든 글자에 «원문 몇 번째 글자?» 꼬리표(offset map) 부착", 10.5, True, INK),
])
arrow(cx, y+0.98, color=C_DET)

# 후보 탐지 3열
y = 3.04
box(LX, y, LW, 1.02, C_DETBG, line=C_DET, line_w=1.4, radius=0.06)
tb(LX+0.12, y+0.03, LW-0.24, 0.26, [("②  후보 탐지  (3종 병렬)  →  candidates[]", 12, True, C_DET)])
cw = (LW-0.48)/3
for i, (t, sub) in enumerate([
    ("M2  regex (9개)", "RRN·전화·카드·이메일·계좌\n+체크섬 검증(Luhn·RRN식)"),
    ("M2.5  dictionary", "은행·기관·학교·병원\n·이름·주소 사전 매칭"),
    ("M5  NER (선택)", "KLUE-RoBERTa v3\n이름·주소·기관 / 기본=Mock"),
]):
    bx = LX+0.16 + i*(cw+0.06)
    box(bx, y+0.30, cw, 0.66, WHITE, line=RGBColor(0xBE,0xD8,0xEE), radius=0.10)
    s1, s2 = sub.split("\n")
    tb(bx, y+0.30, cw, 0.66, [(t, 11.5, True, INK), (s1, 9, False, SUB), (s2, 9, False, SUB)], align=PP_ALIGN.CENTER)
arrow(cx, y+1.04, color=C_REF)

# 보정·정제 3열
y = 4.18
box(LX, y, LW, 1.04, C_REFBG, line=C_REF, line_w=1.4, radius=0.06)
tb(LX+0.12, y+0.03, LW-0.24, 0.26, [("③  보정·정제  (한국어 특화)", 12, True, C_REF)])
for i, (t, sub) in enumerate([
    ("M3  경계보정", "조사/호칭/어미 트림\n«홍길동님께»→«홍길동»"),
    ("M4  문맥가중", "필드라벨«성명:»부스트 / «날씨»\n페널티 / 조합위험→composite승격"),
    ("M6  span 정리", "중복병합·겹침해소(Union-Find)\n주소조각병합·위험등급 상향"),
]):
    bx = LX+0.16 + i*(cw+0.06)
    box(bx, y+0.30, cw, 0.68, WHITE, line=RGBColor(0xB7,0xDD,0xD4), radius=0.10)
    s1, s2 = sub.split("\n")
    tb(bx, y+0.30, cw, 0.68, [(t, 11.5, True, INK), (s1, 8.8, False, SUB), (s2, 8.8, False, SUB)], align=PP_ALIGN.CENTER)
arrow(cx, y+1.06, color=C_POL)

# M7 판정 + 마스킹
y = 5.34
box(LX, y, LW, 0.88, C_POLBG, line=C_POL, line_w=1.4, radius=0.07)
tb(LX+0.16, y+0.07, LW-0.32, 0.38, [
    ("④  M7 PolicyRouter (판정):  위험등급 P0~P3 × 점수 × output_target  →  Action(PASS/MASK/HASH/BLOCK)", 11, True, C_POL)])
tb(LX+0.16, y+0.45, LW-0.32, 0.38, [
    ("    M7 Masker (마스킹):  [REDACTED] / [PERSON_NAME_1] / hmac-sha256:…  치환  ·  BLOCK 하나라도 있으면 전체 차단", 11, False, INK)])
arrow(cx, y+0.90, color=C_AUD)

# M8 감사
y = 6.34
lbox(LX, y, LW, 0.46, C_AUDBG,
     [("⑤  M8  audit_logger  —  HMAC 해시만 기록 · 원문 PII 유출 감지 시 이벤트 폐기 (fail-closed)", 10.5, True, C_AUD)],
     line=C_AUD, line_w=1.4, radius=0.10)
arrow(cx, y+0.48, color=C_IN)

# 출력
y = 6.92
lbox(LX, y, LW, 0.50, C_IN,
     [("출력  GuardrailResponse", 13, True, WHITE),
      ("blocked · masked_text · public_spans(⭐원문 text 필드 없음·value_hash만) · audit_events · metrics", 9.5, False, RGBColor(0xCF,0xDB,0xEC))],
     radius=0.16)

out = r"C:\litellm\프로그램_아키텍처_가운데정렬.pptx"
prs.save(out)
print("SAVED:", out)
