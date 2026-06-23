# -*- coding: utf-8 -*-
"""
한국어 PII 가드레일 — 전체 아키텍처 1-page 슬라이드 생성기
누구나 이해하기 쉽게: 빨강=문제 발견(공격) / 파랑=우리 해결책(방어)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 색 팔레트 ----------
NAVY      = RGBColor(0x16, 0x29, 0x47)   # 타이틀
INK       = RGBColor(0x22, 0x2A, 0x35)   # 본문 진회색
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
PAPER     = RGBColor(0xF4, 0xF6, 0xF9)   # 배경

RED       = RGBColor(0xD6, 0x3A, 0x3A)   # 문제 헤더
RED_BG    = RGBColor(0xFC, 0xEC, 0xEC)   # 문제 패널 배경
RED_BOX   = RGBColor(0xFB, 0xDD, 0xDD)   # 문제 강조 박스

BLUE      = RGBColor(0x1F, 0x77, 0xB4)   # 해결 헤더
TEAL      = RGBColor(0x16, 0x8F, 0x7E)
BLUE_BG   = RGBColor(0xE9, 0xF3, 0xFA)   # 해결 패널 배경
BLUE_BOX  = RGBColor(0xD6, 0xEA, 0xF7)   # 해결 단계 박스

GOLD      = RGBColor(0xE8, 0x9B, 0x12)   # 성과 강조
GREEN     = RGBColor(0x2E, 0x9E, 0x4F)

FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# 배경
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = PAPER
bg.line.fill.background()
bg.shadow.inherit = False


def _set_text(tf, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    """lines: list of (text, size, bold, color, [space_after_pt])"""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(3);  tf.margin_bottom = Pt(3)
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        sp = spec[4] if len(spec) > 4 else 2
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp); p.space_before = Pt(0)
        # 한 paragraph 안에 \n -> run 여러개 대신 단순화: 하나의 run
        r = p.add_run(); r.text = text
        f = r.font
        f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
        # 동아시아 폰트도 동일 지정
        rPr = r._r.get_or_add_rPr()
        ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT}); rPr.append(ea)


def box(x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        line_w=1.0, radius=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    return s


def textbox(x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set_text(tb.text_frame, lines, align, anchor)
    return tb


def labeled_box(x, y, w, h, fill, lines, align=PP_ALIGN.CENTER,
                line=None, line_w=1.0, radius=0.12):
    box(x, y, w, h, fill, line=line, line_w=line_w, radius=radius)
    textbox(x, y, w, h, lines, align=align, anchor=MSO_ANCHOR.MIDDLE)


def arrow_down(cx, y, h=0.28, color=INK):
    w = 0.42
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                               Inches(cx - w/2), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


# ===================== 1) 타이틀 =====================
box(0.3, 0.22, 12.73, 0.92, NAVY, radius=0.14)
textbox(0.55, 0.22, 9.5, 0.92, [
    ("한국어 PII 가드레일 — 전체 구조 한눈에 보기", 27, True, WHITE),
], align=PP_ALIGN.LEFT)
textbox(9.4, 0.22, 3.5, 0.92, [
    ("LiteLLM AI Gateway · v0.2", 12, False, RGBColor(0xC9,0xD6,0xEA)),
    ("개인정보 자동 탐지·차단", 12, False, RGBColor(0xC9,0xD6,0xEA)),
], align=PP_ALIGN.RIGHT)

# ===================== 2) 한 줄 스토리 =====================
box(0.3, 1.28, 12.73, 0.82, WHITE, line=RGBColor(0xD3,0xDA,0xE5), line_w=1.2, radius=0.18)
textbox(0.55, 1.28, 12.2, 0.82, [
    ("핵심 한 줄:  기존 AI 안전장치는 «한국어 개인정보»를 자주 놓친다  →  그래서 우리가 한국어 전용 방패를 만들어 자동으로 가린다.",
     15.5, True, INK),
], align=PP_ALIGN.LEFT)

# ===================== 좌/우 패널 공통 좌표 =====================
PY = 2.30          # 패널 top
PH = 3.95          # 패널 height
LX = 0.30          # 좌 패널 x
LW = 6.05          # 좌 패널 width
RX = 6.98          # 우 패널 x
RW = 6.05          # 우 패널 width

# ----- 좌 패널 배경 (문제) -----
box(LX, PY, LW, PH, RED_BG, line=RGBColor(0xF1,0xC6,0xC6), line_w=1.4, radius=0.06)
# ----- 우 패널 배경 (해결) -----
box(RX, PY, RW, PH, BLUE_BG, line=RGBColor(0xC2,0xDD,0xF0), line_w=1.4, radius=0.06)

# ===================== 좌 패널 내용: STEP 1 문제 증명 =====================
box(LX+0.18, PY+0.16, LW-0.36, 0.62, RED, radius=0.16)
textbox(LX+0.18, PY+0.16, LW-0.36, 0.62, [
    ("①  문제 증명  ·  공격으로 빈틈 찾기", 16.5, True, WHITE),
], align=PP_ALIGN.CENTER)

lcx = LX + LW/2
labeled_box(LX+0.55, PY+0.95, LW-1.10, 0.62, WHITE,
            [("한국어 «함정» 데이터 10,000개 생성", 14, True, INK),
             ("주민번호·전화·진단명 등을 36가지로 변형", 10.5, False, RGBColor(0x6B,0x73,0x80))],
            line=RGBColor(0xE7,0xC7,0xC7))
arrow_down(lcx, PY+1.62, color=RED)

labeled_box(LX+0.55, PY+1.95, LW-1.10, 0.78, WHITE,
            [("기존 5층 안전장치에 그대로 통과시험", 14, True, INK),
             ("Presidio · AWS Bedrock · Lakera · GPT-4o Judge", 10.5, False, RGBColor(0x6B,0x73,0x80))],
            line=RGBColor(0xE7,0xC7,0xC7))
arrow_down(lcx, PY+2.78, color=RED)

labeled_box(LX+0.40, PY+3.10, LW-0.80, 0.68, RED_BOX,
            [("결과: 한국어 텍스트형 개인정보 68% 그대로 통과 ⚠️", 14.5, True, RED)],
            line=RED, line_w=1.4)

# ===================== 우 패널 내용: STEP 2 우리 해결책 =====================
box(RX+0.18, PY+0.16, RW-0.36, 0.62, BLUE, radius=0.16)
textbox(RX+0.18, PY+0.16, RW-0.36, 0.62, [
    ("②  우리 해결책  ·  한국어 전용 방패(v0.2)", 16.5, True, WHITE),
], align=PP_ALIGN.CENTER)

rcx = RX + RW/2
sw = RW - 1.10          # 단계 박스 width
sx = RX + 0.55
step_h = 0.46
gap = 0.16
y = PY + 0.98

steps = [
    ("👤  사용자 입력", "\"제 주민번호 900101-1234567이에요\"", BLUE_BOX, INK),
    ("🔍  탐지", "이름·전화·주민번호·계좌·진단명 자동 찾기", WHITE, INK),
    ("🎭  마스킹", "홍길동 → [이름] / 900101-… → [주민번호]", WHITE, INK),
    ("🤖  AI 안전 호출", "가려진 문장만 GPT에게 전달", WHITE, INK),
    ("✅  안전한 답변", "원본 개인정보는 로그에도 안 남음", BLUE_BOX, INK),
]
for i, (t1, t2, fill, col) in enumerate(steps):
    labeled_box(sx, y, sw, step_h, fill,
                [(f"{t1}    {t2}", 12.5, True if i in (1,2) else False, col)],
                align=PP_ALIGN.LEFT, line=RGBColor(0xBE,0xD8,0xEE))
    # 좌측 라벨 굵게 보이도록: 텍스트는 한 줄로 처리(가시성 우선)
    if i < len(steps) - 1:
        arrow_down(rcx, y + step_h + 0.005, h=gap-0.02, color=BLUE)
    y += step_h + gap

# ===================== 가운데 연결 (발견→해결) =====================
con = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                             Inches(LX+LW-0.02), Inches(PY+PH/2-0.34),
                             Inches(RX-(LX+LW)+0.04), Inches(0.68))
con.fill.solid(); con.fill.fore_color.rgb = GOLD
con.line.fill.background(); con.shadow.inherit = False
textbox(LX+LW-0.05, PY+PH/2-0.34, RX-(LX+LW)+0.10, 0.68,
        [("빈틈\n발견", 9, True, WHITE)], align=PP_ALIGN.CENTER)

# ===================== 3) 성과 바 =====================
by = PY + PH + 0.16
box(0.3, by, 12.73, 0.80, NAVY, radius=0.14)
textbox(0.5, by, 1.7, 0.80, [("📊 성과", 16, True, GOLD)], align=PP_ALIGN.CENTER)

# 성과 3칸
metrics = [
    ("탐지율 80% → 97%", "5층+한국어층 합산"),
    ("한국어 정보유출 26%p 차단", "정규화·사전의 효과"),
    ("11개 게이트웨이 연동", "LiteLLM·Bedrock·LangChain 등"),
]
mx = 2.35
mw = (12.73 + 0.3 - mx - 0.15) / 3
for i, (big, small) in enumerate(metrics):
    cx = mx + i*mw
    textbox(cx, by+0.06, mw-0.15, 0.70, [
        (big, 14.5, True, WHITE),
        (small, 10, False, RGBColor(0xB9,0xC6,0xDB)),
    ], align=PP_ALIGN.CENTER)
    if i < 2:
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(cx+mw-0.10), Inches(by+0.16),
                                    Pt(1.4), Inches(0.48))
        ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x3A,0x4D,0x6E)
        ln.line.fill.background(); ln.shadow.inherit = False

out = r"C:\litellm\한국어_PII_가드레일_아키텍처_1page.pptx"
prs.save(out)
print("SAVED:", out)
