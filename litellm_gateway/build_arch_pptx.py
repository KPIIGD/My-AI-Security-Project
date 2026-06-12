"""
한국어 PII 가드레일 — 아키텍처 & 학습 전략 PPT 생성
7장 구성: 표지 / 전체 아키텍처 / L0 정규화 / 4 Detector / Span Ledger / NER 학습 전략 / 데이터 흐름
"""
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ── 디자인 시스템 ──
FONT = "맑은 고딕"
NAVY = RGBColor(0x1E, 0x3A, 0x8A)        # 제목·강조
DARK = RGBColor(0x1F, 0x29, 0x37)        # 본문 텍스트
GRAY = RGBColor(0x6B, 0x72, 0x80)        # 보조 텍스트
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)       # 박스 배경
BORDER = RGBColor(0xCB, 0xD5, 0xE1)      # 테두리
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
RED = RGBColor(0xDC, 0x26, 0x26)         # 강조 (수치)
GREEN = RGBColor(0x05, 0x96, 0x69)       # 성공
ORANGE = RGBColor(0xEA, 0x58, 0x0C)      # 경고/주의

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]  # 보통 6번이 blank, 없으면 0번 폴백


def get_blank():
    # 안전하게 빈 레이아웃 찾기
    for lo in prs.slide_layouts:
        if "blank" in (lo.name or "").lower():
            return lo
    return prs.slide_layouts[-1]


BLANK = get_blank()


def set_run(run, text, size=14, bold=False, color=DARK, font=FONT):
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color, font=font)
    return tb


def add_lines(slide, x, y, w, h, lines, size=12, color=DARK, line_spacing=1.15,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(item, tuple):
            txt, kwargs = item
        else:
            txt, kwargs = item, {}
        set_run(p.add_run(), txt,
                size=kwargs.get("size", size),
                bold=kwargs.get("bold", False),
                color=kwargs.get("color", color),
                font=kwargs.get("font", FONT))
    return tb


def add_box(slide, x, y, w, h, fill=LIGHT, line=BORDER, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    # 그림자 제거 (XML 직접)
    sppr = shp._element.spPr
    eff = sppr.find(qn("a:effectLst"))
    if eff is not None:
        sppr.remove(eff)
    sppr.append(etree.fromstring(f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))
    shp.text_frame.text = ""
    shp.text_frame.margin_left = Emu(80000)
    shp.text_frame.margin_right = Emu(80000)
    shp.text_frame.margin_top = Emu(60000)
    shp.text_frame.margin_bottom = Emu(60000)
    return shp


def box_with_title_lines(slide, x, y, w, h, title, lines, fill=LIGHT,
                         title_color=NAVY, title_size=14, line_size=11,
                         line_color=DARK):
    """라운드 박스 + 제목 + 본문 라인. 한 번에."""
    add_box(slide, x, y, w, h, fill=fill)
    # 제목
    add_text(slide, x, y, w, Inches(0.45), title,
             size=title_size, bold=True, color=title_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 본문
    body_y = y + Inches(0.5)
    body_h = h - Inches(0.55)
    add_lines(slide, x, body_y, w, body_h, lines, size=line_size,
              color=line_color, align=PP_ALIGN.LEFT)


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, w=1.5):
    line = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(w)
    # 화살표 끝
    ln = line.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return line


def slide_title(slide, num_text, title):
    """슬라이드 상단 헤더 — 번호 + 제목 + 가로선"""
    add_text(slide, Inches(0.5), Inches(0.3), Inches(0.8), Inches(0.5),
             num_text, size=12, bold=True, color=GRAY)
    add_text(slide, Inches(1.3), Inches(0.3), Inches(11), Inches(0.5),
             title, size=24, bold=True, color=NAVY)
    # 가로선
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(0.92),
                                      Inches(12.83), Inches(0.92))
    line.line.color.rgb = NAVY
    line.line.width = Pt(2)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 1 — 표지
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# 좌상단 작은 라벨
add_text(s, Inches(0.7), Inches(0.7), Inches(5), Inches(0.5),
         "캡스톤 시스템 설계 · 2026.05", size=12, bold=True, color=GRAY)

# 중앙 메인 타이틀
add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(1.0),
         "한국어 PII 가드레일", size=44, bold=True, color=NAVY,
         align=PP_ALIGN.LEFT)

add_text(s, Inches(0.7), Inches(3.3), Inches(12), Inches(0.7),
         "시스템 아키텍처 & 학습 전략", size=26, bold=False, color=DARK,
         align=PP_ALIGN.LEFT)

# 부제 (영문)
add_text(s, Inches(0.7), Inches(4.1), Inches(12), Inches(0.5),
         "Korean PII Guardrail — L0 결정론 + KLUE-RoBERTa NER + Tier 4 문맥 룰",
         size=14, color=GRAY)

# 하단 라인 + 정보
line = s.shapes.add_connector(1, Inches(0.7), Inches(6.3),
                              Inches(12.6), Inches(6.3))
line.line.color.rgb = NAVY
line.line.width = Pt(1.5)

add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.4),
         "정보보안학과 캡스톤 · 지도교수 임정묵", size=12, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 2 — 전체 아키텍처
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "01", "전체 아키텍처")

# 입력 (좌상단 4개 박스)
y0 = Inches(1.2)
input_w = Inches(2.4)
input_h = Inches(0.7)
gap = Inches(0.2)

inputs = [
    ("📝 텍스트", "Plain"),
    ("📄 PDF", "PyMuPDF + pdfplumber"),
    ("📊 XLSX", "openpyxl 행 concat"),
    ("📋 CSV", "pandas"),
]
ix = Inches(0.7)
for title, desc in inputs:
    box_with_title_lines(s, ix, y0, input_w, input_h, title,
                        [(desc, {"size": 9, "color": GRAY})],
                        title_size=12)
    ix += input_w + gap

# 입력 → 통합 텍스트 화살표
add_arrow(s, Inches(6.66), y0 + input_h, Inches(6.66), y0 + input_h + Inches(0.35))

# 통합 텍스트 박스
add_box(s, Inches(3.5), Inches(2.25), Inches(6.33), Inches(0.5),
        fill=RGBColor(0xE0, 0xE7, 0xFF))
add_text(s, Inches(3.5), Inches(2.25), Inches(6.33), Inches(0.5),
         "통합 텍스트 + offset_meta(page, row, col)", size=12, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 화살표
add_arrow(s, Inches(6.66), Inches(2.78), Inches(6.66), Inches(3.13))

# L0 정규화 박스
box_with_title_lines(s, Inches(2.5), Inches(3.15), Inches(8.33), Inches(0.65),
                    "🧹 L0 한국어 정규화 13단계",
                    [("자모분해 · 동형문자 · 야민정음 · 공백 등 결정론 전처리",
                      {"size": 10, "color": GRAY})],
                    fill=RGBColor(0xFE, 0xF3, 0xC7))

# 화살표
add_arrow(s, Inches(6.66), Inches(3.83), Inches(6.66), Inches(4.18))

# 4 Detector 병렬
y_det = Inches(4.2)
det_w = Inches(2.95)
det_h = Inches(1.45)
detectors = [
    ("① 정규식 42→60", "RRN/PHONE/CARD\nEMAIL/JWT/IP", RGBColor(0xDB, 0xEA, 0xFE)),
    ("② 사전 22→27", "병명/처방/회사\n학교/직책 외", RGBColor(0xD1, 0xFA, 0xE5)),
    ("③ NER (inference)", "klue/roberta-base\nNAME/ADDR/ORG", RGBColor(0xFE, 0xE2, 0xE2)),
    ("④ Tier 4 룰", "3-slot quasi-id\nrisk_flag", RGBColor(0xFE, 0xF3, 0xC7)),
]
dx = Inches(0.7)
for title, desc, fill in detectors:
    add_box(s, dx, y_det, det_w, det_h, fill=fill)
    add_text(s, dx, y_det + Inches(0.1), det_w, Inches(0.4),
             title, size=13, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, dx, y_det + Inches(0.55), det_w, Inches(0.85),
             desc, size=10, color=DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    dx += det_w + Inches(0.1)

# 화살표 4개 → Span Ledger
y_arr_start = y_det + det_h
y_arr_end = y_arr_start + Inches(0.35)
for i in range(4):
    cx = Inches(0.7) + det_w / 2 + i * (det_w + Inches(0.1))
    add_arrow(s, cx, y_arr_start, Inches(6.66), y_arr_end)

# Span Ledger
box_with_title_lines(s, Inches(3.0), Inches(6.05), Inches(7.33), Inches(0.65),
                    "📚 Span Ledger — 우선순위·merge·tie-break",
                    [("정규식 > NER > 사전 > Tier 4  ·  span 충돌 시 confidence 비교",
                      {"size": 10, "color": GRAY})],
                    fill=RGBColor(0xE0, 0xE7, 0xFF))

# 출력
add_arrow(s, Inches(6.66), Inches(6.72), Inches(6.66), Inches(7.0))
add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.4),
         "📤 비식별 텍스트 + risk_flag + span 목록(offset 보존)",
         size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 3 — L0 정규화 13단계
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "02", "L0 한국어 정규화 13단계")

steps = [
    ("①", "NFKC", "유니코드 정규화"),
    ("②", "ZWSP / ZWNJ", "Zero-Width Space 제거"),
    ("③", "Soft Hyphen", "U+00AD 제거"),
    ("④", "Combining Mark", "결합 문자 제거"),
    ("⑤", "Fullwidth", "전각 → 반각 ASCII"),
    ("⑥", "Mathematical Bold", "𝟬𝟭 → 0,1"),
    ("⑦", "Circled Digits", "①② → 1,2"),
    ("⑧", "한자 성씨", "金 → 김"),
    ("⑨", "★ 자모 결합", "ㅈㅜㅁㅣㄴ → 주민"),
    ("⑩", "Homoglyph", "키릴 о → 라틴 o"),
    ("⑪", "공백", "정규화"),
    ("⑫", "제어 문자", "제거"),
    ("⑬", "Trim", "최종 정리"),
]

# 4×4 그리드 (13개)
cols = 4
rows = 4
gx = Inches(0.7)
gy = Inches(1.3)
cw = (Inches(12.6) - Inches(0.7)) / cols - Inches(0.15)
ch = Inches(1.15)

for i, (num, name, desc) in enumerate(steps):
    r = i // cols
    c = i % cols
    bx = gx + c * (cw + Inches(0.15))
    by = gy + r * (ch + Inches(0.15))
    fill = RGBColor(0xFE, 0xF3, 0xC7) if "★" in num else LIGHT
    add_box(s, bx, by, cw, ch, fill=fill)
    add_text(s, bx, by + Inches(0.08), cw, Inches(0.35),
             num, size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bx, by + Inches(0.42), cw, Inches(0.35),
             name, size=12, bold=True, color=DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bx, by + Inches(0.78), cw, Inches(0.32),
             desc, size=9, color=GRAY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# 하단 메모
add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
         "★ 자모 결합 단계가 KoJailFuzz 적대적 변이의 핵심 방어선",
         size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 4 — 4 Detector 책임 분담
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "03", "4 Detector 책임 분담")

# 4개 박스 큰 사이즈로
det_y = Inches(1.3)
det_w = Inches(2.95)
det_h = Inches(4.5)
det_x = Inches(0.7)

detectors_full = [
    {
        "title": "① 정규식",
        "subtitle": "42 → 60개",
        "tier": "Tier 1 정형 PII",
        "color": RGBColor(0xDB, 0xEA, 0xFE),
        "items": [
            "주민등록번호",
            "전화번호",
            "이메일",
            "신용카드",
            "계좌번호",
            "JWT 토큰",
            "세션 토큰",
            "여권/면허",
            "차량번호",
            "GPS 좌표",
            "AWS 키",
        ],
        "method": "결정론적 패턴 매칭",
    },
    {
        "title": "② 사전 기반",
        "subtitle": "22 → 27 카테고리",
        "tier": "Tier 2/3 텍스트형",
        "color": RGBColor(0xD1, 0xFA, 0xE5),
        "items": [
            "알레르기 / 진단명",
            "처방약 / 수술 이력",
            "혈액형 / 장애",
            "종교 / 혼인",
            "성별 / 성지향",
            "학교 / 학위",
            "회사 / 부서 / 직책",
            "+ salary / debt",
            "+ family / location",
            "+ job-title-honorific",
        ],
        "method": "키워드 + 컨텍스트 매칭",
    },
    {
        "title": "③ NER",
        "subtitle": "inference-only",
        "tier": "Tier 3 문맥형 PII",
        "color": RGBColor(0xFE, 0xE2, 0xE2),
        "items": [
            "NAME (이름)",
            "ADDRESS (주소)",
            "ORG (조직)",
            "TITLE (직책)",
            "DEPT (부서)",
            "—",
            "klue/roberta-base",
            "사전학습 체크포인트",
            "(full fine-tune 미진행)",
            "DISEASE/TITLE/DEPT는",
            "사전 기반으로 분리",
        ],
        "method": "통계적 (BPE 토크나이저)",
    },
    {
        "title": "④ Tier 4 룰",
        "subtitle": "재식별 위험 플래그",
        "tier": "Tier 4 quasi-id",
        "color": RGBColor(0xFE, 0xF3, 0xC7),
        "items": [
            "3-slot 조합 발화 시 flag:",
            "—",
            "  가족 / 본인",
            "  +",
            "  질환 / 처방",
            "  +",
            "  지역 / 나이 / 직장",
            "—",
            "출력: risk_flag 만",
            "(mask 금지)",
            "FPR<5% 검증 필수",
        ],
        "method": "결정론 룰 (조합형)",
    },
]

for i, d in enumerate(detectors_full):
    bx = det_x + i * (det_w + Inches(0.13))
    add_box(s, bx, det_y, det_w, det_h, fill=d["color"])
    # 제목
    add_text(s, bx, det_y + Inches(0.1), det_w, Inches(0.4),
             d["title"], size=16, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bx, det_y + Inches(0.5), det_w, Inches(0.3),
             d["subtitle"], size=11, color=GRAY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bx, det_y + Inches(0.85), det_w, Inches(0.3),
             d["tier"], size=10, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 항목
    item_lines = [(line, {"size": 10}) for line in d["items"]]
    add_lines(s, bx, det_y + Inches(1.25), det_w, Inches(2.85),
              item_lines, size=10, line_spacing=1.15,
              align=PP_ALIGN.CENTER)
    # 방식 (하단)
    add_text(s, bx, det_y + Inches(4.05), det_w, Inches(0.4),
             d["method"], size=9, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 하단 한 줄
add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.4),
         "→ 4 detector 결과를 Span Ledger가 통합 (다음 슬라이드)",
         size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 5 — Span Ledger
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "04", "Span Ledger — 통합·중재 메커니즘")

# 좌측: 입력 + 4 detector 출력
y_top = Inches(1.3)
add_text(s, Inches(0.7), y_top, Inches(5.5), Inches(0.4),
         "입력 텍스트:", size=12, bold=True, color=NAVY)
add_box(s, Inches(0.7), y_top + Inches(0.4), Inches(5.5), Inches(0.6),
        fill=RGBColor(0xF8, 0xFA, 0xFC))
add_text(s, Inches(0.7), y_top + Inches(0.4), Inches(5.5), Inches(0.6),
         '"홍길동 010-1234-5678 페니실린 알레르기"',
         size=11, color=DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         font="Consolas")

# 4 detector 출력 메타
y_meta = Inches(2.2)
add_text(s, Inches(0.7), y_meta, Inches(5.5), Inches(0.4),
         "4 Detector 출력 메타 (treat as ledger entries):",
         size=12, bold=True, color=NAVY)

ledger_items = [
    ("정규식", "PHONE", "(5, 18)", "0.99", RGBColor(0xDB, 0xEA, 0xFE)),
    ("NER", "NAME", "(0, 3)", "0.92", RGBColor(0xFE, 0xE2, 0xE2)),
    ("사전", "ALLERGY", "(19, 25)", "0.95", RGBColor(0xD1, 0xFA, 0xE5)),
    ("Tier 4", "—", "—", "no flag", RGBColor(0xFE, 0xF3, 0xC7)),
]
yy = Inches(2.65)
for src, ptype, span, conf, fill in ledger_items:
    add_box(s, Inches(0.7), yy, Inches(5.5), Inches(0.45), fill=fill)
    add_text(s, Inches(0.85), yy, Inches(1.4), Inches(0.45),
             src, size=11, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.3), yy, Inches(1.5), Inches(0.45),
             ptype, size=11, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.7), yy, Inches(1.5), Inches(0.45),
             span, size=10, color=GRAY, font="Consolas",
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.1), yy, Inches(1.0), Inches(0.45),
             conf, size=10, bold=True, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    yy += Inches(0.5)

# 우측: 룰 박스
y_right = Inches(1.3)
right_x = Inches(7.2)
right_w = Inches(5.6)
add_box(s, right_x, y_right, right_w, Inches(5.3), fill=LIGHT)
add_text(s, right_x, y_right + Inches(0.1), right_w, Inches(0.4),
         "Span Ledger 통합 규칙",
         size=15, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

rules = [
    ("🔑 핵심 원칙", {"size": 12, "bold": True, "color": NAVY}),
    ("• 원문 보존 — placeholder 치환 X", {"size": 11}),
    ("• BPE OOV 토큰 노이즈 회피", {"size": 11}),
    ("• train-test mismatch 차단", {"size": 11}),
    (" ", {"size": 6}),
    ("📋 우선순위 (충돌 시 적용)", {"size": 12, "bold": True, "color": NAVY}),
    ("정규식  >  NER  >  사전  >  Tier 4 룰",
     {"size": 12, "bold": True, "color": DARK, "font": "Consolas"}),
    (" ", {"size": 6}),
    ("⚖️ Tie-break", {"size": 12, "bold": True, "color": NAVY}),
    ("• Span 겹치면 confidence 높은 쪽 채택", {"size": 11}),
    ("• Confidence 동률 → detector 우선순위", {"size": 11}),
    ("• offset_meta(page, row, col) 보존", {"size": 11}),
]
add_lines(s, right_x + Inches(0.3), y_right + Inches(0.55),
          right_w - Inches(0.6), Inches(4.6),
          rules, size=11, color=DARK, line_spacing=1.3)

# 최종 출력
add_arrow(s, Inches(3.45), Inches(5.2), Inches(3.45), Inches(5.65), color=GREEN, w=2)
add_box(s, Inches(0.7), Inches(5.7), Inches(5.5), Inches(1.2),
        fill=RGBColor(0xDC, 0xFC, 0xE7))
add_text(s, Inches(0.7), Inches(5.75), Inches(5.5), Inches(0.4),
         "최종 출력",
         size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(6.15), Inches(5.5), Inches(0.7),
         '"●●● 010-●●●●-●●●● 페니실린 ●●●●●"\n+ risk_flag: false',
         size=11, color=DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         font="Consolas")


# ═══════════════════════════════════════════════════════════════════════════
# Slide 6 — NER 학습 전략
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "05", "NER 학습 전략 — Inference-only 보완 모듈")

# 좌측: full fine-tune 비교 (X)
fx = Inches(0.7)
fy = Inches(1.3)
fw = Inches(5.8)
fh = Inches(2.5)
add_box(s, fx, fy, fw, fh, fill=RGBColor(0xFE, 0xE2, 0xE2))
add_text(s, fx, fy + Inches(0.1), fw, Inches(0.4),
         "❌ Full Fine-tune (탈락)",
         size=15, bold=True, color=RED,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_lines(s, fx + Inches(0.3), fy + Inches(0.55), fw - Inches(0.6), Inches(1.9),
          [
              ("• 데이터 라벨링 1주", {"size": 11}),
              ("• 학습 1주 + 오탐 튜닝 1주", {"size": 11}),
              ("• 합계 3주+ → 6주 일정 무리", {"size": 11}),
              (" ", {"size": 6}),
              ("• 합성 데이터 품질 리스크 큼", {"size": 11}),
              ("• 클래스 불균형 (RRN 1% vs NAME 40%)", {"size": 11}),
              ("• 외부 평가 시 in-distribution 의심", {"size": 11}),
          ],
          size=11, color=DARK, line_spacing=1.25)

# 우측: inference-only (O)
ix = Inches(6.83)
add_box(s, ix, fy, fw, fh, fill=RGBColor(0xDC, 0xFC, 0xE7))
add_text(s, ix, fy + Inches(0.1), fw, Inches(0.4),
         "✅ Inference-only (채택)",
         size=15, bold=True, color=GREEN,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_lines(s, ix + Inches(0.3), fy + Inches(0.55), fw - Inches(0.6), Inches(1.9),
          [
              ("• klue/roberta-base 사전학습 그대로", {"size": 11}),
              ("• 학습 시간 0", {"size": 11}),
              ("• 통합·데모·오류분석 산출물 확보", {"size": 11}),
              (" ", {"size": 6}),
              ("• KLUE-NER 라벨 → PII 라벨 매핑 룰", {"size": 11}),
              ("• 후처리로 NAME/ADDR/ORG/TITLE/DEPT", {"size": 11}),
              ("• DISEASE 등은 사전 기반 분리", {"size": 11}),
          ],
          size=11, color=DARK, line_spacing=1.25)

# 하단: 매핑 룰 다이어그램
my = Inches(4.0)
add_text(s, Inches(0.7), my, Inches(12), Inches(0.4),
         "KLUE-NER 라벨 → PII 라벨 매핑 흐름",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# 박스 5개: KLUE-RoBERTa → KLUE-NER 라벨 → 매핑 → PII 라벨 → Span
flow_y = Inches(4.55)
flow_h = Inches(1.1)
boxes = [
    ("klue/roberta-base", "사전학습 체크포인트", RGBColor(0xE0, 0xE7, 0xFF)),
    ("KLUE-NER 라벨", "PER · LOC · ORG", RGBColor(0xFE, 0xE2, 0xE2)),
    ("매핑 룰", "PER → NAME\nLOC → ADDRESS\nORG → ORG", RGBColor(0xFE, 0xF3, 0xC7)),
    ("PII 라벨", "NAME · ADDRESS · ORG\n+ TITLE/DEPT/DISEASE\n(사전 분리)", RGBColor(0xD1, 0xFA, 0xE5)),
    ("Span 출력", "(start, end, type,\n confidence)", RGBColor(0xE0, 0xE7, 0xFF)),
]
bw = Inches(2.3)
gap_b = Inches(0.15)
total_w = bw * 5 + gap_b * 4
bx = (Inches(13.333) - total_w) / 2
for title, desc, color in boxes:
    add_box(s, bx, flow_y, bw, flow_h, fill=color)
    add_text(s, bx, flow_y + Inches(0.05), bw, Inches(0.35),
             title, size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bx, flow_y + Inches(0.4), bw, Inches(0.65),
             desc, size=9, color=DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    bx += bw + gap_b

# 화살표 4개
ax_y = flow_y + flow_h / 2
for i in range(4):
    cx_start = (Inches(13.333) - total_w) / 2 + bw * (i + 1) + gap_b * i
    cx_end = cx_start + gap_b
    add_arrow(s, cx_start, ax_y, cx_end, ax_y, color=GRAY, w=1.5)

# 하단 메모
add_text(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.4),
         "→ 학습 0주, 라벨 매핑 룰 작성 1주, 오류분석 1주 (총 2주, W3 완료)",
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 7 — 데이터 흐름 + 학습 파이프라인
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "06", "데이터 흐름 & 학습 파이프라인")

# 좌측: 데이터 확보 (4 stage)
ly = Inches(1.3)
add_text(s, Inches(0.7), ly, Inches(5.8), Inches(0.4),
         "데이터 확보 (Stage 0~4)", size=15, bold=True, color=NAVY)

stages = [
    ("Stage 0  ★ W1 첫날 선행", "라벨 가이드 1페이지\nNAME / ADDR 경계 정의\n→ BIO 태그 5~7개 통폐합",
     RGBColor(0xFE, 0xF3, 0xC7)),
    ("Stage 1  W1", "Faker-ko 합성 1만 건\n+ KoJailFuzz augmentation 20~30%\n(문서 단위 8:1:1 분할 후)",
     RGBColor(0xDB, 0xEA, 0xFE)),
    ("Stage 2  W2", "KLUE-NER 재라벨링\nPER → NAME, LOC → ADDR\nDISEASE/TITLE 사전 분리",
     RGBColor(0xD1, 0xFA, 0xE5)),
    ("Stage 4  W2 (필수)", "외부 평가셋 준비\n• KLUE-NER test split\n• KoJailFuzz 14종",
     RGBColor(0xFE, 0xE2, 0xE2)),
]
sy = Inches(1.75)
sh_h = Inches(1.15)
for title, desc, color in stages:
    add_box(s, Inches(0.7), sy, Inches(5.8), sh_h, fill=color)
    add_text(s, Inches(0.85), sy + Inches(0.08), Inches(5.5), Inches(0.35),
             title, size=12, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.85), sy + Inches(0.42), Inches(5.5), Inches(0.7),
             desc, size=10, color=DARK,
             anchor=MSO_ANCHOR.TOP)
    sy += sh_h + Inches(0.1)

# 우측: 학습 파이프라인
rx = Inches(7.0)
rw = Inches(5.7)
add_text(s, rx, ly, rw, Inches(0.4),
         "학습 파이프라인 (W1~W6)", size=15, bold=True, color=NAVY)

py_ = Inches(1.75)
phases = [
    ("W1", "라벨 가이드 + Faker-ko 합성 + 8:1:1 분할", RGBColor(0xFE, 0xF3, 0xC7)),
    ("W2", "KLUE-NER 재라벨링 + 외부 평가셋", RGBColor(0xDB, 0xEA, 0xFE)),
    ("W3", "★ NER inference + Tier 4 룰 작성", RGBColor(0xFE, 0xE2, 0xE2)),
    ("W4", "Span Ledger 통합 + L0 강화 (정규식 60, 사전 27)", RGBColor(0xD1, 0xFA, 0xE5)),
    ("W5", "PDF/XLSX/CSV 어댑터", RGBColor(0xE0, 0xE7, 0xFF)),
    ("W6", "평가 (span-F1 / KoJailFuzz / latency / FPR) + 발표", RGBColor(0xFE, 0xF3, 0xC7)),
]
for week, task, color in phases:
    add_box(s, rx, py_, rw, Inches(0.65), fill=color)
    add_text(s, rx + Inches(0.15), py_, Inches(0.7), Inches(0.65),
             week, size=14, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, rx + Inches(0.95), py_, rw - Inches(1.0), Inches(0.65),
             task, size=11, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    py_ += Inches(0.72)

# 하단: 평가 메트릭 4종
mx = Inches(0.7)
my = Inches(6.5)
mw = Inches(12.6)
add_box(s, mx, my, mw, Inches(0.85), fill=LIGHT)
add_text(s, mx, my + Inches(0.05), mw, Inches(0.3),
         "평가 메트릭 (4종)", size=11, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)
add_text(s, mx, my + Inches(0.35), mw, Inches(0.45),
         "① span-exact entity-F1 per category    "
         "② KoJailFuzz 14종 F1 유지율    "
         "③ latency p50/p95 (T4)    "
         "④ parse coverage % (문서, PII F1과 분리)",
         size=10, color=DARK, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════════════════════════
import os
DST = r"c:\litellm\한국어PII가드레일_아키텍처_학습전략.pptx"
prs.save(DST)
print(f"[OK] 슬라이드: {len(prs.slides)}")
print(f"[OK] 파일: {DST}")
print(f"[OK] 크기: {os.path.getsize(DST)/1024:.1f} KB")
