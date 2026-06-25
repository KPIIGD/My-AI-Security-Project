# -*- coding: utf-8 -*-
"""
베이스라인(L1~L3) 한국어 PII 우회 — 두 평가본 나란히 비교 슬라이드
좌: run_b_10k_summary.json (payloads_10k)  /  우: phase2_v4_baseline.json (orig)
모든 수치는 git 정본(My-AI-Security-Project) 기준.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY   = RGBColor(0x16, 0x29, 0x47)
INK    = RGBColor(0x22, 0x2A, 0x35)
SUB    = RGBColor(0x55, 0x61, 0x72)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PAPER  = RGBColor(0xF4, 0xF6, 0xF9)
ROW_A  = RGBColor(0xFF, 0xFF, 0xFF)
ROW_B  = RGBColor(0xED, 0xF1, 0xF6)
HEADC  = RGBColor(0x2B, 0x3C, 0x57)
WARN_BG= RGBColor(0xFB, 0xDF, 0xDF)
WARN   = RGBColor(0xCB, 0x33, 0x33)

GREEN_D= RGBColor(0x16, 0x7A, 0x3C)
GREEN  = RGBColor(0x2E, 0x9E, 0x4F)
OLIVE  = RGBColor(0xB5, 0x7A, 0x0E)
RED    = RGBColor(0xCB, 0x33, 0x33)

FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333); prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = PAPER; bg.line.fill.background(); bg.shadow.inherit = False


def _set(tf, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        sp = spec[4] if len(spec) > 4 else 1
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp); p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
        rPr = r._r.get_or_add_rPr(); rPr.append(rPr.makeelement(qn('a:ea'), {'typeface': FONT}))


def box(x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0, radius=0.08):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: s.adjustments[0] = radius
        except Exception: pass
    return s


def tb(x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set(t.text_frame, lines, align, anchor); return t


def true_color(v):
    if v >= 95: return GREEN_D
    if v >= 78: return GREEN
    if v >= 60: return OLIVE
    return RED


def draw_table(x0, y0, title_main, title_sub, rows):
    W = 6.0
    cw = [2.62, 1.74, 1.64]   # 구분 / TRUE / 우회
    rh = 0.475
    # 제목
    box(x0, y0, W, 0.56, HEADC, radius=0.10, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(x0+0.12, y0, W-0.24, 0.56, [(title_main, 12.5, True, WHITE), (title_sub, 9.5, False, RGBColor(0xB9,0xC6,0xDB))])
    # 헤더행
    hy = y0 + 0.62
    heads = ["구분", "TRUE 탐지율", "우회"]
    cx = x0
    for j, htext in enumerate(heads):
        box(cx, hy, cw[j], rh, HEADC)
        tb(cx, hy, cw[j], rh, [(htext, 11, True, WHITE)],
           align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
        cx += cw[j]
    # 데이터행
    for i, (label, tv, bv, indent, warn) in enumerate(rows):
        ry = hy + rh*(i+1)
        base = WARN_BG if warn else (ROW_B if i % 2 else ROW_A)
        cx = x0
        # 구분 셀
        box(cx, ry, cw[0], rh, base)
        tb(cx + (0.28 if indent else 0.0), ry, cw[0]-(0.30 if indent else 0.04), rh,
           [(("└ " if indent else "") + label, 11 if not warn else 11.5, (not indent) or warn,
             WARN if warn else INK)])
        cx += cw[0]
        # TRUE 셀
        box(cx, ry, cw[1], rh, base)
        tb(cx, ry, cw[1], rh, [(f"{tv:.2f}%", 12.5 if warn else 12, True, true_color(tv))], align=PP_ALIGN.CENTER)
        cx += cw[1]
        # 우회 셀
        box(cx, ry, cw[2], rh, base)
        warn_mark = " ⚠️" if warn else ""
        tb(cx, ry, cw[2], rh, [(f"{bv:.2f}%{warn_mark}", 12 if warn else 11.5, warn,
                                WARN if warn else SUB)], align=PP_ALIGN.CENTER)
    return 0.62 + rh*(len(rows)+1)

# ===== 타이틀 =====
box(0.3, 0.20, 12.73, 0.74, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
tb(0.55, 0.20, 12.2, 0.74, [("베이스라인(L1~L3, Layer 0 OFF) 한국어 PII 우회 — 두 평가본 교차 비교", 21, True, WHITE)])

# ===== 핵심 한 줄 =====
box(0.3, 1.04, 12.73, 0.62, WHITE, line=RGBColor(0xD3,0xDA,0xE5), line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
tb(0.5, 1.04, 12.33, 0.62, [
    ("측정본이 달라도 결론은 같다 →  영어 PII 99%대 탐지,  한국어 «텍스트형(semantic)»은 절반(48~52%)이 그냥 우회.",
     13.5, True, INK)])

# ===== 좌 표: run_b =====
rows_b = [
    ("전체",            79.01, 20.99, False, False),
    ("영어(EN)",        99.23,  0.77, False, False),
    ("한국어(KR)",      68.19, 31.81, False, False),
    ("KR 체크섬",       81.61, 18.39, True,  False),
    ("KR 정규식형",     72.32, 27.68, True,  False),
    ("KR 텍스트형(semantic)", 47.93, 52.07, True, True),
]
draw_table(0.40, 1.86, "① run_b  (전체 79.01%)", "출처: PII/results/summaries/run_b_10k_summary.json", rows_b)

# ===== 우 표: phase2_v4 =====
rows_v4 = [
    ("전체",            80.15, 19.85, False, False),
    ("영어(EN)",        99.37,  0.63, False, False),
    ("한국어(KR)",      69.86, 30.14, False, False),
    ("KR 체크섬",       83.14, 16.86, True,  False),
    ("KR 정규식형",     74.00, 26.00, True,  False),
    ("KR 텍스트형(semantic)", 49.62, 50.38, True, True),
]
draw_table(6.93, 1.86, "② phase2_v4  (전체 80.15%)", "출처: PII/evaluation/phase2_v4_baseline.json (orig)", rows_v4)

# ===== 하단 결론 =====
ey = 5.86
box(0.3, ey, 12.73, 1.32, RGBColor(0xEC,0xF3, 0xEC), line=GREEN, line_w=1.6, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
tb(0.5, ey+0.10, 12.33, 0.40, [("✅  두 측정본이 독립적으로 같은 결론 → 발견은 «측정 오차»가 아니라 «구조적 약점»", 14, True, GREEN_D)])
tb(0.5, ey+0.52, 12.33, 0.72, [
    ("•  영어 PII 99%대 vs 한국어 텍스트형 ~48% — 같은 가드레일이 언어·유형에 따라 2배 격차.   "
     "•  두 표 차이(79.01 vs 80.15)는 «동일 1만건»을 집계 기준만 미세조정한 차이(같은 raw: eval_10k_l1l3.json).", 10.8, False, INK, 3),
    ("•  체크섬·정규식형(숫자 PII)은 한국어도 72~83% 잡힘 → 진짜 구멍은 «텍스트형»(진단·알레르기·직책 등). 이게 Layer 0의 핵심 타깃.", 10.8, False, INK, 0),
])

# 각주
tb(0.3, 7.22, 12.73, 0.24, [
    ("데이터 정본: github.com/KPIIGD(또는 vmaca123)/My-AI-Security-Project — PII/results/, PII/evaluation/  ·  C:\\litellm 사본은 git 미추적",
     8.5, False, RGBColor(0x8A,0x93,0xA1))], align=PP_ALIGN.CENTER)

out = r"C:\litellm\베이스라인_두평가본_비교.pptx"
prs.save(out)
print("SAVED:", out)
