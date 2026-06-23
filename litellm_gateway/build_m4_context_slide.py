# -*- coding: utf-8 -*-
"""
M4 문맥가중(Context Scoring) 발표 전용 슬라이드 1장
같은 글자도 '주변 단어'에 따라 PII가 되기도/안 되기도 한다.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY   = RGBColor(0x16, 0x29, 0x47)
INK    = RGBColor(0x22, 0x2A, 0x35)
SUB    = RGBColor(0x55, 0x61, 0x72)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PAPER  = RGBColor(0xF4, 0xF6, 0xF9)

GREEN  = RGBColor(0x1E, 0x8E, 0x4F); GREEN_BG = RGBColor(0xE7, 0xF5, 0xEC)
RED    = RGBColor(0xCB, 0x33, 0x33); RED_BG   = RGBColor(0xFC, 0xEC, 0xEC)
PURPLE = RGBColor(0x6E, 0x3A, 0xA6); PUR_BG   = RGBColor(0xF1, 0xEA, 0xF8)
GOLD   = RGBColor(0xB5, 0x7A, 0x0E); GOLD_BG  = RGBColor(0xFB, 0xF3, 0xDD)

FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333); prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = PAPER; bg.line.fill.background(); bg.shadow.inherit = False


def _set(tf, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(7); tf.margin_right = Pt(7); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        sp = spec[4] if len(spec) > 4 else 1
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp); p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
        rPr = r._r.get_or_add_rPr(); rPr.append(rPr.makeelement(qn('a:ea'), {'typeface': FONT}))


def box(x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=1.0, radius=0.08):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: s.adjustments[0] = radius
        except Exception: pass
    return s


def tb(x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set(t.text_frame, lines, align, anchor); return t


# ===== 타이틀 =====
box(0.3, 0.20, 12.73, 0.80, NAVY, radius=0.14)
tb(0.55, 0.20, 12.2, 0.80, [("M4  문맥가중 — 같은 글자도 «주변 단어»가 PII인지 가른다", 23, True, WHITE)])

# ===== 핵심 한 줄 =====
box(0.3, 1.18, 12.73, 0.66, WHITE, line=RGBColor(0xD3,0xDA,0xE5), line_w=1.2, radius=0.16)
tb(0.5, 1.18, 12.33, 0.66, [
    ("정규식·사전은 «글자»만 본다.  M4는 «문장 속 주변 단어»를 읽어  ① 점수를 올리고  ② 내리고  ③ 둘 이상 모이면 위험등급을 올린다.  (LLM 없이 규칙으로)",
     13, True, INK)], align=PP_ALIGN.LEFT)

# ===== 3 카드 =====
cy, ch = 1.98, 2.74
gap = 0.22
cw = (12.73 - 2*gap) / 3
xs = [0.3, 0.3+cw+gap, 0.3+2*(cw+gap)]

# ① 부스트
box(xs[0], cy, cw, ch, GREEN_BG, line=GREEN, line_w=1.6, radius=0.05)
box(xs[0], cy, cw, 0.50, GREEN, radius=0.05)
tb(xs[0], cy, cw, 0.50, [("①  부스트  ▲ 점수 올림", 13.5, True, WHITE)], align=PP_ALIGN.CENTER)
tb(xs[0]+0.14, cy+0.52, cw-0.28, 0.34, [("«진짜 PII» 단서가 있으면 +", 10.5, True, GREEN)])
tb(xs[0]+0.14, cy+0.86, cw-0.28, ch-1.0, [
    ("•  «성명·이름·고객명» 라벨 → 이름  +0.25", 10, False, INK, 4),
    ("•  «계좌·입금·송금» 라벨 → 계좌  +0.30", 10, False, INK, 4),
    ("•  은행명 동반(우리은행…) → 계좌  +0.25", 10, False, INK, 4),
    ("•  이름 뒤 «님·씨·과장·교수» → +0.15", 10, False, INK, 4),
    ("•  같은 문장에 전화 동반 → 이름 +0.20", 10, False, INK, 4),
    ("•  «소속·근무·재직» → 조직  +0.25", 10, False, INK, 4),
], anchor=MSO_ANCHOR.TOP)

# ② 페널티
box(xs[1], cy, cw, ch, RED_BG, line=RED, line_w=1.6, radius=0.05)
box(xs[1], cy, cw, 0.50, RED, radius=0.05)
tb(xs[1], cy, cw, 0.50, [("②  페널티  ▼ 점수 내림", 13.5, True, WHITE)], align=PP_ALIGN.CENTER)
tb(xs[1]+0.14, cy+0.52, cw-0.28, 0.34, [("«PII 아님» 단서가 있으면 −  (오탐 방지)", 10.5, True, RED)])
tb(xs[1]+0.14, cy+0.86, cw-0.28, ch-1.0, [
    ("•  «날씨·맑네요·비·눈» → 이름  −0.35", 10, False, INK, 4),
    ("•  «대표번호·고객센터·콜센터» → −0.25", 10, False, INK, 4),
    ("•  «상호·식당·카페·브랜드» → 이름 −0.25", 10, False, INK, 4),
    ("•  «stack·json·로그·error» → −0.20", 10, False, INK, 4),
    ("•  «예시·샘플·테스트·더미» → −0.15", 10, False, INK, 4),
    ("•  사설 IP(10.· 192.168.…) → −0.25", 10, False, INK, 4),
], anchor=MSO_ANCHOR.TOP)

# ③ 조합 승격
box(xs[2], cy, cw, ch, PUR_BG, line=PURPLE, line_w=1.6, radius=0.05)
box(xs[2], cy, cw, 0.50, PURPLE, radius=0.05)
tb(xs[2], cy, cw, 0.50, [("③  조합 승격  ⊕ 위험등급↑", 13.5, True, WHITE)], align=PP_ALIGN.CENTER)
tb(xs[2]+0.14, cy+0.52, cw-0.28, 0.34, [("한 문장에 PII가 둘 이상 → 개인 특정 가능", 10.5, True, PURPLE)])
tb(xs[2]+0.14, cy+0.86, cw-0.28, ch-1.0, [
    ("•  이름 + 전화  →  P1", 10, False, INK, 4),
    ("•  이름 + 이메일 / 주소 / 계좌  →  P1", 10, False, INK, 4),
    ("•  고객ID + 전화  →  P1", 10, False, INK, 4),
    ("•  환자번호 + 병원  →  P1", 10, False, INK, 4),
    ("•  생일 + 호수 + 학교  →  P1", 10, False, INK, 4),
    ("이름만 = 약함,  이름+전화 = 특정 가능 → 위험↑", 9.3, False, SUB, 4),
], anchor=MSO_ANCHOR.TOP)

# ===== 하단: 대표 예시 문장 =====
ey = 4.96
box(0.3, ey, 12.73, 2.28, GOLD_BG, line=GOLD, line_w=1.6, radius=0.04)
tb(0.5, ey+0.10, 12.33, 0.36, [("⭐ 한 문장으로 보는 M4  —  «같은 김하늘, 문맥이 다르면 결과도 다르다»", 14, True, GOLD)])

# 예시 A (부스트+조합)
box(0.55, ey+0.55, 8.0, 1.55, WHITE, line=GREEN, line_w=1.4, radius=0.05)
tb(0.7, ey+0.62, 7.7, 0.40, [("✅  «고객명 김하늘 님, 연락처 010-1234-5678입니다»", 13.5, True, INK)])
tb(0.7, ey+1.02, 7.7, 1.0, [
    ("«고객명» → 김하늘 이름  ▲+0.25      «님» → 호칭  ▲+0.15", 10.5, False, GREEN, 3),
    ("«연락처» → 전화  ▲+0.15      이름 + 전화 같은 문장  ⊕ 조합 P1 승격", 10.5, False, GREEN, 3),
    ("⟹  둘 다 «확실한 PII»로 판정 → 강하게 마스킹", 11, True, INK, 0),
], anchor=MSO_ANCHOR.TOP)

# 예시 B (페널티)
box(8.72, ey+0.55, 4.28, 1.55, WHITE, line=RED, line_w=1.4, radius=0.05)
tb(8.87, ey+0.62, 3.98, 0.40, [("❌  «오늘 하늘이 참 맑네요»", 13.5, True, INK)])
tb(8.87, ey+1.02, 3.98, 1.0, [
    ("«맑네요·날씨» 단서 감지  ▼−0.35", 10.5, False, RED, 3),
    ("⟹  여기서 «하늘»은 이름이 아님", 10.5, True, INK, 3),
    ("→ 마스킹 안 함 (오탐 방지)", 10.5, False, SUB, 0),
], anchor=MSO_ANCHOR.TOP)

out = r"C:\litellm\M4_문맥가중_발표슬라이드.pptx"
prs.save(out)
print("SAVED:", out)
